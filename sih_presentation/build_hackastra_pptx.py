import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_hackastra_verbatim_deck(output_pptx_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    cur_dir = os.path.dirname(os.path.abspath(__file__))
    logo_png = os.path.join(cur_dir, "sih_logo.png")
    bulb_png = os.path.join(cur_dir, "sih_bulb.png")
    img_dir = os.path.join(cur_dir, "images")
    fra_ui = os.path.join(img_dir, "fra_setu_ui.png")
    fc_img = os.path.join(img_dir, "hackastra_flowchart.png")
    arch_img = os.path.join(img_dir, "hackastra_arch.png")

    NAVY_TITLE = RGBColor(27, 54, 93)        # #1B365D
    BLACK = RGBColor(0, 0, 0)
    BLUE_HEADING = RGBColor(0, 112, 192)     # #0070C0
    GREEN_ACCENT = RGBColor(21, 128, 61)     # #15803D
    WHITE = RGBColor(255, 255, 255)

    def add_common_header(slide, title_text, slide_num, team_name="Hackastra"):
        # 1. Top Left Team Oval
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.45), Inches(0.18), Inches(1.85), Inches(0.75))
        oval.fill.solid()
        oval.fill.fore_color.rgb = WHITE
        oval.line.color.rgb = BLACK
        oval.line.width = Pt(1.5)
        tf_oval = oval.text_frame
        tf_oval.word_wrap = True
        p1 = tf_oval.paragraphs[0]
        p1.text = "Team"
        p1.font.name = "Arial"
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = BLACK
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf_oval.add_paragraph()
        p2.text = team_name
        p2.font.name = "Arial"
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = BLACK
        p2.alignment = PP_ALIGN.CENTER

        # 2. Center Pill Title
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.6), Inches(0.18), Inches(6.1), Inches(0.75))
        pill.fill.solid()
        pill.fill.fore_color.rgb = WHITE
        pill.line.color.rgb = BLACK
        pill.line.width = Pt(1.5)
        tf_pill = pill.text_frame
        tf_pill.word_wrap = True
        p_p = tf_pill.paragraphs[0]
        p_p.text = title_text
        p_p.font.name = "Arial"
        p_p.font.size = Pt(24)
        p_p.font.bold = True
        p_p.font.color.rgb = BLACK
        p_p.alignment = PP_ALIGN.CENTER

        # 3. Top Right SIH Logo
        if os.path.exists(logo_png):
            slide.shapes.add_picture(logo_png, Inches(11.0), Inches(0.15), width=Inches(1.85))

        # 4. Slide Number at Bottom Right
        num_box = slide.shapes.add_textbox(Inches(12.4), Inches(6.92), Inches(0.6), Inches(0.4))
        p_num = num_box.text_frame.paragraphs[0]
        p_num.text = str(slide_num)
        p_num.font.name = "Arial"
        p_num.font.size = Pt(14)
        p_num.font.bold = True
        p_num.font.color.rgb = BLACK
        p_num.alignment = PP_ALIGN.RIGHT

    def add_sharp_border(slide, x, y, w, h):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BLACK
        box.line.width = Pt(1.8)
        return box

    # =========================================================================
    # SLIDE 1: TITLE PAGE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)

    # Top Title Centered
    h1_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.32), Inches(12.333), Inches(0.8))
    p_h1 = h1_box.text_frame.paragraphs[0]
    p_h1.text = "SMART INDIA HACKATHON 2025"
    p_h1.font.name = "Times New Roman"
    p_h1.font.size = Pt(36)
    p_h1.font.bold = True
    p_h1.font.color.rgb = NAVY_TITLE
    p_h1.alignment = PP_ALIGN.CENTER

    # Top Right SIH Logo
    if os.path.exists(logo_png):
        slide1.shapes.add_picture(logo_png, Inches(10.8), Inches(0.2), width=Inches(2.0))

    # Center Title: TITLE PAGE
    tp_box = slide1.shapes.add_textbox(Inches(1.5), Inches(1.3), Inches(10.333), Inches(0.6))
    p_tp = tp_box.text_frame.paragraphs[0]
    p_tp.text = "TITLE PAGE"
    p_tp.font.name = "Arial"
    p_tp.font.size = Pt(30)
    p_tp.font.bold = True
    p_tp.font.color.rgb = BLACK
    p_tp.alignment = PP_ALIGN.CENTER

    # Right SIH Brain Bulb Watermark (Centered background)
    if os.path.exists(bulb_png):
        slide1.shapes.add_picture(bulb_png, Inches(5.2), Inches(1.5), width=Inches(3.8))

    # Left Metadata Content
    meta_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.15), Inches(12.0), Inches(4.8))
    tf_m = meta_box.text_frame
    tf_m.word_wrap = True

    items = [
        ("• Problem Statement ID –", "SIH25108", True),
        ("• Problem Statement Title-", "Development of AI-powered FRA Atlas and WebGIS-based Decision Support System (DSS) for Integrated Monitoring of Forest Rights Act (FRA) Implementation.\n(States to be concentrated: Madhya Pradesh, Tripura , Odisha, Telangana)", False),
        ("• Theme-", "Miscellaneous", True),
        ("• PS Category- ", "Software", True),
        ("• Team ID- ", "57385", False),
        ("• Team Name- ", "Team Hackastra", False)
    ]

    for idx, (label, val, val_under) in enumerate(items):
        p = tf_m.paragraphs[0] if idx == 0 else tf_m.add_paragraph()
        p.space_after = Pt(13)
        r_lbl = p.add_run()
        r_lbl.text = label
        r_lbl.font.name = "Arial"
        r_lbl.font.size = Pt(17)
        r_lbl.font.bold = True
        r_lbl.font.color.rgb = BLACK

        r_val = p.add_run()
        r_val.text = val
        r_val.font.name = "Arial"
        r_val.font.size = Pt(16.5) if "Title" in label else Pt(17)
        r_val.font.bold = True
        r_val.font.underline = val_under
        r_val.font.color.rgb = BLACK

    # Slide 1 Number
    num1_box = slide1.shapes.add_textbox(Inches(12.4), Inches(6.92), Inches(0.6), Inches(0.4))
    p1_num = num1_box.text_frame.paragraphs[0]
    p1_num.text = "1"
    p1_num.font.name = "Arial"
    p1_num.font.size = Pt(14)
    p1_num.font.bold = True
    p1_num.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 2: IDEA TITLE
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_common_header(slide2, "IDEA TITLE", 2, "Hackastra")

    # Left Sharp Box
    add_sharp_border(slide2, 0.45, 1.08, 6.0, 5.75)
    s2_left = slide2.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(5.8), Inches(5.6))
    tf2_l = s2_left.text_frame
    tf2_l.word_wrap = True

    p_s2_head = tf2_l.paragraphs[0]
    p_s2_head.text = "❖ Proposed Solution:"
    p_s2_head.font.name = "Arial"
    p_s2_head.font.size = Pt(18)
    p_s2_head.font.bold = True
    p_s2_head.font.color.rgb = BLUE_HEADING
    p_s2_head.space_after = Pt(6)

    sol_bullets = [
        ("NGO Digital Aid :", "Teach Gram Sabhas and communities apps, IVR, and GIS tools."),
        ("Digital Witness & Testimony:", "Voice and video records stored securely as valid proof."),
        ("Title Insurance / Digital Patta :", "Ensures authenticity with secure 7/12 land record entry."),
        ("Grievance Resolution & Illegal Alerts:", "Efficient FRA grievance handling and unauthorized land encroachment warnings."),
        ("AI Disciplinary Dashboard :", "Suggests actions and tracks officer accountability."),
        ("Analytics-Driven Officer Monitoring:", "Tracks verification time, claim outcomes, and grievances reported."),
        ("Interactive FRA Web-GIS Atlas:", "Dynamic platform for exploring forest rights and assets."),
        ("AI Disputes Mediator:", "Resolving conflicts within existing legal frameworks."),
        ("GI Tagging and E-Signature :", "To preserve cultural heritage safeguard traditional products and paperless approval .")
    ]

    for title, desc in sol_bullets:
        p = tf2_l.add_paragraph()
        p.space_after = Pt(4)
        r_t = p.add_run()
        r_t.text = f"• {title}\n"
        r_t.font.name = "Arial"
        r_t.font.size = Pt(11)
        r_t.font.bold = True
        r_t.font.underline = True
        r_t.font.color.rgb = GREEN_ACCENT

        r_d = p.add_run()
        r_d.text = f"  {desc}"
        r_d.font.name = "Arial"
        r_d.font.size = Pt(10.5)
        r_d.font.color.rgb = BLACK

    # Right Sharp Box with Flowchart Image
    add_sharp_border(slide2, 6.65, 1.08, 6.25, 5.75)
    if os.path.exists(fc_img):
        slide2.shapes.add_picture(fc_img, Inches(6.7), Inches(1.13), width=Inches(6.15), height=Inches(5.65))

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_common_header(slide3, "TECHNICAL APPROACH", 3, "Hackastra")

    # Left Top Box: 10 Bullets
    add_sharp_border(slide3, 0.45, 1.08, 6.0, 2.85)
    s3_top = slide3.shapes.add_textbox(Inches(0.55), Inches(1.12), Inches(5.8), Inches(2.75))
    tf3_t = s3_top.text_frame
    tf3_t.word_wrap = True

    tech_bullets = [
        ("Data Extraction:", "Optical Character Recogntion(OCR), NLP."),
        ("WebGIS platform:", "Leaflet.js(for interactive maps)."),
        ("Asset Mapping:", "Random Forest & CNN."),
        ("Data Storage:", "PostgreSQl & PostGIS(for geospatial data)."),
        ("Data Management:", "Cloud Object Storage & MySQL(Structured)."),
        ("GeoServer:", "Use standard protocols WMS and WFS to serve maps."),
        ("Decision Support System:", "for scheme recommendations."),
        ("Blockchain based Security:", "Hyperledger Fabric."),
        ("Key Lock Framework:", "Open-source Identity & Access Management"),
        ("Rasa:", "Preferred for its offline, multilingual, and secure chatbot.")
    ]

    for idx, (title, desc) in enumerate(tech_bullets):
        p = tf3_t.paragraphs[0] if idx == 0 else tf3_t.add_paragraph()
        p.space_after = Pt(2)
        r_t = p.add_run()
        r_t.text = f"• {title} "
        r_t.font.name = "Arial"
        r_t.font.size = Pt(9.5)
        r_t.font.bold = True
        r_t.font.underline = True
        r_t.font.color.rgb = GREEN_ACCENT

        r_d = p.add_run()
        r_d.text = desc
        r_d.font.name = "Arial"
        r_d.font.size = Pt(9.5)
        r_d.font.color.rgb = BLACK

    # Left Bottom Box: UI Screenshot
    add_sharp_border(slide3, 0.45, 4.05, 6.0, 2.78)
    if os.path.exists(fra_ui):
        slide3.shapes.add_picture(fra_ui, Inches(0.48), Inches(4.08), width=Inches(5.94), height=Inches(2.45))

    cap_box = slide3.shapes.add_textbox(Inches(0.5), Inches(6.55), Inches(5.9), Inches(0.25))
    p_c = cap_box.text_frame.paragraphs[0]
    p_c.text = "User Interface: It is a draft version, subject to future updates."
    p_c.font.name = "Arial"
    p_c.font.size = Pt(9.5)
    p_c.font.bold = True
    p_c.font.color.rgb = BLACK

    # Right Sharp Box: Architecture Image
    add_sharp_border(slide3, 6.65, 1.08, 6.25, 5.75)
    if os.path.exists(arch_img):
        slide3.shapes.add_picture(arch_img, Inches(6.7), Inches(1.13), width=Inches(6.15), height=Inches(5.65))

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_common_header(slide4, "FEASIBILITY AND VIABILITY", 4, "Hackastra")

    # Left Box: Feasibility
    add_sharp_border(slide4, 0.45, 1.08, 6.0, 5.75)
    s4_left = slide4.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(5.8), Inches(5.6))
    tf4_l = s4_left.text_frame
    tf4_l.word_wrap = True

    p_f = tf4_l.paragraphs[0]
    p_f.text = "❖ Feasibility:"
    p_f.font.name = "Arial"
    p_f.font.size = Pt(18)
    p_f.font.bold = True
    p_f.font.color.rgb = BLUE_HEADING
    p_f.space_after = Pt(4)

    feas_tree = [
        ("1. Technical Feasibility:", [
            ("1.Digitization, AI Mapping and Web-GIS Dashboards:", [
                "OCR and AI extract data from documents.",
                "Web-GIS dashboards show real-time interactive maps."
            ]),
            ("2. DSS Integration with Blockchain Security:", [
                "AI engines match schemes to community needs.",
                "Blockchain ensures secure, tamper-proof digital records."
            ])
        ]),
        ("2. Operational Feasibility:", [
            ("1.Existing Setup and Simple Staff Training:", [
                "Government data exists, needs proper connection.",
                "Staff need basic training through short workshops."
            ]),
            ("2. Stakeholder Support and Effective FRA Implementation:", [
                "Ministries and NGOs already support FRA schemes.",
                "System strengthens efforts, not replacing existing work."
            ])
        ]),
        ("3. Legal & Ethical Feasibility:", [
            ("1.Data Privacy:", [
                "Must follow India’s DPDP Act, ensuring legal compliance.",
                "Blockchain and encryption safeguard sensitive user data."
            ]),
            ("2. Consent Mechanism:", [
                "FRA holders’ informed consent required in all processes.",
                "Real-time feedback enables secure two-way communication."
            ])
        ])
    ]

    for sec_title, sub_list in feas_tree:
        p_sec = tf4_l.add_paragraph()
        p_sec.space_after = Pt(2)
        r_s = p_sec.add_run()
        r_s.text = sec_title
        r_s.font.name = "Arial"
        r_s.font.size = Pt(11)
        r_s.font.bold = True
        r_s.font.underline = True
        r_s.font.color.rgb = GREEN_ACCENT

        for sub_head, b_list in sub_list:
            p_sub = tf4_l.add_paragraph()
            p_sub.space_after = Pt(1)
            r_sh = p_sub.add_run()
            r_sh.text = sub_head
            r_sh.font.name = "Arial"
            r_sh.font.size = Pt(10.5)
            r_sh.font.bold = True
            r_sh.font.color.rgb = BLACK

            for b in b_list:
                p_b = tf4_l.add_paragraph()
                p_b.space_after = Pt(1)
                r_b = p_b.add_run()
                r_b.text = f"  • {b}"
                r_b.font.name = "Arial"
                r_b.font.size = Pt(10)
                r_b.font.color.rgb = BLACK

    # Right Box: Viability
    add_sharp_border(slide4, 6.65, 1.08, 6.25, 5.75)
    s4_right = slide4.shapes.add_textbox(Inches(6.75), Inches(1.15), Inches(6.05), Inches(5.6))
    tf4_r = s4_right.text_frame
    tf4_r.word_wrap = True

    p_v = tf4_r.paragraphs[0]
    p_v.text = "❖ Viability:"
    p_v.font.name = "Arial"
    p_v.font.size = Pt(18)
    p_v.font.bold = True
    p_v.font.color.rgb = BLUE_HEADING
    p_v.space_after = Pt(4)

    viab_tree = [
        ("1. Economic Viability:", [
            ("1.Cost-Effective & Efficient:", [
                "Saves paperwork, effort, duplication, and reduces fraud.",
                "Cuts processing delays, expenses, and operational overhead."
            ]),
            ("2. Supported & Scalable Technology:", [
                "Backed by government programs and tribal inclusion schemes.",
                "Open-source tech ensures flexible, low-cost,future-ready deployment."
            ])
        ]),
        ("2. Social Viability:", [
            ("1. Empowers Communities & Improves Scheme Access:", [
                "Provides transparency, secure records, better documentation.",
                "Directly benefits FRA holders through scheme access."
            ]),
            ("2. Reduces Conflict & Promotes Stronger Inclusion:", [
                "Verified records, alerts reduce frequent land disputes.",
                "Mobile updates, two-way communication improve engagement."
            ])
        ]),
        ("3. Scalability:", [
            ("1. Pan-India Expansion for Wider Community Benefits:", [
                "Pilot in districts, then scale across India.",
                "Adjust platform easily for local data."
            ]),
            ("2. Cross-Sector Use of WebGIS + DSS:", [
                "Adaptable for agriculture, urban development needs.",
                "Useful in disaster management and planning."
            ])
        ])
    ]

    for sec_title, sub_list in viab_tree:
        p_sec = tf4_r.add_paragraph()
        p_sec.space_after = Pt(2)
        r_s = p_sec.add_run()
        r_s.text = sec_title
        r_s.font.name = "Arial"
        r_s.font.size = Pt(11)
        r_s.font.bold = True
        r_s.font.underline = True
        r_s.font.color.rgb = GREEN_ACCENT

        for sub_head, b_list in sub_list:
            p_sub = tf4_r.add_paragraph()
            p_sub.space_after = Pt(1)
            r_sh = p_sub.add_run()
            r_sh.text = sub_head
            r_sh.font.name = "Arial"
            r_sh.font.size = Pt(10.5)
            r_sh.font.bold = True
            r_sh.font.color.rgb = BLACK

            for b in b_list:
                p_b = tf4_r.add_paragraph()
                p_b.space_after = Pt(1)
                r_b = p_b.add_run()
                r_b.text = f"  • {b}"
                r_b.font.name = "Arial"
                r_b.font.size = Pt(10)
                r_b.font.color.rgb = BLACK

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_common_header(slide5, "IMPACT AND BENEFITS", 5, "Hackastra")

    # Left Box: Benefits
    add_sharp_border(slide5, 0.45, 1.08, 6.0, 5.75)
    s5_left = slide5.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(5.8), Inches(5.6))
    tf5_l = s5_left.text_frame
    tf5_l.word_wrap = True

    p_ben = tf5_l.paragraphs[0]
    p_ben.text = "❖ Benefits:"
    p_ben.font.name = "Arial"
    p_ben.font.size = Pt(18)
    p_ben.font.bold = True
    p_ben.font.color.rgb = BLUE_HEADING
    p_ben.space_after = Pt(4)

    ben_tree = [
        ("1.Social Benefits:", [
            "Digital platform enables communities with land rights tracking.",
            "IVR helpline, e-learning ensure inclusive literacy access.",
            "Transparent records, video testimonies reduce conflict, build trust."
        ]),
        ("2. Environmental Benefits:", [
            "Sustainable forest use by recognized communities.",
            "Conservation through community stewardship.",
            "GI Tagging preserves cultural heritage & traditional products."
        ]),
        ("3. Economic Benefits:", [
            "Access to schemes (MGNREGA, PM-KISAN, DAGAJU, Jal Jeevan Mission)",
            "GI Tagging boosts local products & tribal livelihoods",
            "Title insurance & digital patta ensure secure land ownership → easier loans & income stability"
        ]),
        ("4. Governance Benefits:", [
            "Audit trails ensure accountability, prevent misuse.",
            "Whistleblowing portal boosts transparency, fights corruption.",
            "Cross-validation stops duplication and fraudulent claims."
        ]),
        ("5.Legal & Security Benefits:", [
            "E-Signature & Digital Patta provide legally valid ownership records.",
            "Protection against land encroachment (Illegal Practice Alert).",
            "Digital witness/testimony ensures fair dispute settlement."
        ])
    ]

    for sec_title, b_list in ben_tree:
        p_sec = tf5_l.add_paragraph()
        p_sec.space_after = Pt(2)
        r_s = p_sec.add_run()
        r_s.text = sec_title
        r_s.font.name = "Arial"
        r_s.font.size = Pt(11)
        r_s.font.bold = True
        r_s.font.underline = True
        r_s.font.color.rgb = GREEN_ACCENT

        for b in b_list:
            p_b = tf5_l.add_paragraph()
            p_b.space_after = Pt(1.5)
            r_b = p_b.add_run()
            r_b.text = f"  • {b}"
            r_b.font.name = "Arial"
            r_b.font.size = Pt(10)
            r_b.font.color.rgb = BLACK

    # Right Box: Impacts
    add_sharp_border(slide5, 6.65, 1.08, 6.25, 5.75)
    s5_right = slide5.shapes.add_textbox(Inches(6.75), Inches(1.15), Inches(6.05), Inches(5.6))
    tf5_r = s5_right.text_frame
    tf5_r.word_wrap = True

    p_imp = tf5_r.paragraphs[0]
    p_imp.text = "❖ Impacts:"
    p_imp.font.name = "Arial"
    p_imp.font.size = Pt(18)
    p_imp.font.bold = True
    p_imp.font.color.rgb = BLUE_HEADING
    p_imp.space_after = Pt(4)

    imp_tree = [
        ("1.Secure Land Rights & Transparency:", [
            "Communities track claims, reducing corruption and fraud.",
            "Digital patta ensures authenticity in land ownership records.",
            "Builds trust between citizens and government authorities."
        ]),
        ("2. Forest Protection & Conservation with AI:", [
            "AI monitors land to prevent illegal encroachment.",
            "Conserves biodiversity and safeguards natural resources.",
            "Promotes sustainable, community-driven forest management."
        ]),
        ("3. Boost Livelihoods & Local Development:", [
            "Resource planning improves incomes and rural economy.",
            "GI tagging protects heritage, boosts local products.",
            "FRA enables fair access to welfare schemes."
        ]),
        ("4. Digital, Eco-Friendly & Cost Effective:", [
            "Reduces paperwork,ensuring greener, eco-friendly governance.",
            "Preserves tribal knowledge through digital archiving tools.",
            "Lowers administrative costs and increases efficiency."
        ]),
        ("5. Smart, Transparent & Scalable Governance:", [
            "Transparent dashboards track officer performance and claims.",
            "AI dispute mediation ensures fairness, reduces human bias.",
            "Unified digital platform enables future-ready expansion."
        ])
    ]

    for sec_title, b_list in imp_tree:
        p_sec = tf5_r.add_paragraph()
        p_sec.space_after = Pt(2)
        r_s = p_sec.add_run()
        r_s.text = sec_title
        r_s.font.name = "Arial"
        r_s.font.size = Pt(11)
        r_s.font.bold = True
        r_s.font.underline = True
        r_s.font.color.rgb = GREEN_ACCENT

        for b in b_list:
            p_b = tf5_r.add_paragraph()
            p_b.space_after = Pt(1.5)
            r_b = p_b.add_run()
            r_b.text = f"  • {b}"
            r_b.font.name = "Arial"
            r_b.font.size = Pt(10)
            r_b.font.color.rgb = BLACK

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_common_header(slide6, "RESEARCH AND REFERENCES", 6, "Hackastra")

    add_sharp_border(slide6, 0.45, 1.08, 12.45, 5.75)
    s6_box = slide6.shapes.add_textbox(Inches(0.6), Inches(1.18), Inches(12.15), Inches(5.5))
    tf6 = s6_box.text_frame
    tf6.word_wrap = True

    refs = [
        ("1. FRA: ", "The Forest Rights Act (FRA), 2006 gives rights to forest-dwelling and tribal people, but challenges in implementation remain.\nhttps://tribal.nic.in/FRA.aspx\nhttps://mpvanmitra.mkcl.org/hi"),
        ("2. Best Practices for Web Security: ", "This guide outlines key web security measures to safeguard user information using Blockchain Technology.\nhttps://owasp.org/www-project-blockchain-appsec-standard/"),
        ("3. The Wildlife Protection Act: ", "1972 aims to protect wild animals, birds, plants, and their habitats.\nhttps://share.google/F3sPygTbDmIslAmAd")
    ]

    for idx, (head, body) in enumerate(refs):
        p = tf6.paragraphs[0] if idx == 0 else tf6.add_paragraph()
        p.space_after = Pt(10)
        r_h = p.add_run()
        r_h.text = head
        r_h.font.name = "Arial"
        r_h.font.size = Pt(11.5)
        r_h.font.bold = True
        r_h.font.underline = True
        r_h.font.color.rgb = GREEN_ACCENT

        r_b = p.add_run()
        r_b.text = body
        r_b.font.name = "Arial"
        r_b.font.size = Pt(11)
        r_b.font.color.rgb = BLACK

    p_lead = tf6.add_paragraph()
    p_lead.space_after = Pt(8)
    r_l = p_lead.add_run()
    r_l.text = "To enable smooth implementation and integration of the FRA Atlas System, the following contacts serve as official state-level resources for communication, data coordination, and support:"
    r_l.font.name = "Arial"
    r_l.font.size = Pt(11)
    r_l.font.color.rgb = BLACK

    contacts = [
        ("4. Madhya Pradesh: ", "dirtadp@mp.gov.in"),
        ("5. Odisha: ", "https://stsc.odisha.gov.in/acts-policies-guidelines/scheduled-tribes-and-other-traditional-forest-dwellers"),
        ("6. Tripura: ", "https://forest.tripura.gov.in/"),
        ("7. Telangana: ", "https://forestrights.telangana.gov.in/")
    ]

    for head, val in contacts:
        p_c = tf6.add_paragraph()
        p_c.space_after = Pt(5)
        r_h = p_c.add_run()
        r_h.text = head
        r_h.font.name = "Arial"
        r_h.font.size = Pt(11)
        r_h.font.bold = True
        r_h.font.underline = True
        r_h.font.color.rgb = GREEN_ACCENT

        r_v = p_c.add_run()
        r_v.text = val
        r_v.font.name = "Arial"
        r_v.font.size = Pt(11)
        r_v.font.color.rgb = BLUE_HEADING

    prs.save(output_pptx_path)
    print(f"Generated 100% Exact Hackastra PPTX at: {output_pptx_path}")

if __name__ == "__main__":
    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "SIH_Hackastra_Exact_Submission.pptx")
    build_hackastra_verbatim_deck(out_path)
