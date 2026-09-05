import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation(output_pptx_path):
    prs = Presentation()
    # 16:9 Widescreen standard (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    cur_dir = os.path.dirname(os.path.abspath(__file__))
    logo_png = os.path.join(cur_dir, "sih_logo.png")
    bulb_png = os.path.join(cur_dir, "sih_bulb.png")
    vx_logo_png = os.path.join(cur_dir, "virasatx_logo.png")
    
    img_dir = os.path.join(cur_dir, "images")
    home_img = os.path.join(img_dir, "home.png")
    map_img = os.path.join(img_dir, "map.png")
    ai_img = os.path.join(img_dir, "ai_guide.png")
    discover_img = os.path.join(img_dir, "discover.png")
    traditions_img = os.path.join(img_dir, "living_traditions.png")

    # Precise Color Palette Matching Reference 1 & 2
    NAVY = RGBColor(0, 43, 73)            # #002B49
    ROYAL_BLUE = RGBColor(0, 112, 192)    # #0070C0
    FOOTER_BLUE = RGBColor(0, 90, 156)    # #005A9C
    DARK_GREEN = RGBColor(21, 128, 61)    # #15803D
    ORANGE = RGBColor(234, 88, 12)        # #EA580C
    BLACK = RGBColor(0, 0, 0)
    DARK_GREY = RGBColor(51, 65, 85)      # #334155
    SLATE = RGBColor(71, 85, 105)         # #475569
    LIGHT_BG = RGBColor(248, 250, 252)    # #F8FAFC
    BORDER_GREY = RGBColor(203, 213, 225) # #CBD5E1
    WHITE = RGBColor(255, 255, 255)
    GOLD = RGBColor(147, 107, 56)         # #936B38

    # Helper: Add Slide Header
    def add_header(slide, title_text, subtitle_text=None, is_pill=False, innovation_star=False):
        # 1. Left Team Logo + Branding
        if os.path.exists(vx_logo_png):
            slide.shapes.add_picture(vx_logo_png, Inches(0.4), Inches(0.2), width=Inches(0.55), height=Inches(0.55))
        
        tx_team = slide.shapes.add_textbox(Inches(1.05), Inches(0.18), Inches(2.8), Inches(0.6))
        tf_t = tx_team.text_frame
        tf_t.word_wrap = True
        p1 = tf_t.paragraphs[0]
        p1.text = "VirasatX"
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(17)
        p1.font.bold = True
        p1.font.color.rgb = NAVY
        
        p2 = tf_t.add_paragraph()
        p2.text = "India's Heritage Repository"
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(8.5)
        p2.font.bold = True
        p2.font.color.rgb = ROYAL_BLUE

        # 2. Center Title
        if is_pill:
            pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.6), Inches(0.18), Inches(6.1), Inches(0.62))
            pill.fill.solid()
            pill.fill.fore_color.rgb = WHITE
            pill.line.color.rgb = BLACK
            pill.line.width = Pt(1.5)
            tf_p = pill.text_frame
            tf_p.word_wrap = True
            p = tf_p.paragraphs[0]
            p.text = title_text
            p.font.name = "Segoe UI"
            p.font.size = Pt(19)
            p.font.bold = True
            p.font.color.rgb = BLACK
            p.alignment = PP_ALIGN.CENTER
            if subtitle_text:
                p_sub = tf_p.add_paragraph()
                p_sub.text = subtitle_text
                p_sub.font.name = "Segoe UI"
                p_sub.font.size = Pt(9.5)
                p_sub.font.color.rgb = DARK_GREY
                p_sub.alignment = PP_ALIGN.CENTER
        else:
            tx_title = slide.shapes.add_textbox(Inches(3.4), Inches(0.15), Inches(6.5), Inches(0.65))
            tf_t = tx_title.text_frame
            tf_t.word_wrap = True
            p = tf_t.paragraphs[0]
            p.text = title_text
            p.font.name = "Segoe UI"
            p.font.size = Pt(21)
            p.font.bold = True
            p.font.color.rgb = NAVY if "SMART INDIA" in title_text else BLACK
            p.alignment = PP_ALIGN.CENTER

        if innovation_star:
            tx_star = slide.shapes.add_textbox(Inches(9.8), Inches(0.22), Inches(1.5), Inches(0.4))
            p_s = tx_star.text_frame.paragraphs[0]
            p_s.text = "★ - Innovation"
            p_s.font.name = "Segoe UI"
            p_s.font.size = Pt(11)
            p_s.font.bold = True
            p_s.font.color.rgb = RGBColor(217, 119, 6)

        # 3. Right SIH Logo
        if os.path.exists(logo_png):
            slide.shapes.add_picture(logo_png, Inches(11.2), Inches(0.15), width=Inches(1.75))

    # Helper: Add Slide Footer
    def add_footer(slide, slide_num):
        footer_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35))
        footer_bar.fill.solid()
        footer_bar.fill.fore_color.rgb = FOOTER_BLUE
        footer_bar.line.fill.background()

        # Text handle left
        tx_f = slide.shapes.add_textbox(Inches(0.4), Inches(7.15), Inches(4.0), Inches(0.35))
        tf_f = tx_f.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = "@ V I R A S A T X"
        p_f.font.name = "Segoe UI"
        p_f.font.size = Pt(12)
        p_f.font.bold = True
        p_f.font.color.rgb = WHITE

        # Slide Number right
        tx_num = slide.shapes.add_textbox(Inches(12.2), Inches(7.15), Inches(0.7), Inches(0.35))
        tf_num = tx_num.text_frame
        p_num = tf_num.paragraphs[0]
        p_num.text = str(slide_num)
        p_num.font.name = "Segoe UI"
        p_num.font.size = Pt(13)
        p_num.font.bold = True
        p_num.font.color.rgb = WHITE
        p_num.alignment = PP_ALIGN.RIGHT

    # =========================================================================
    # SLIDE 1: TITLE PAGE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    add_header(slide1, "SMART INDIA HACKATHON 2026")
    add_footer(slide1, 1)

    # Title Page Subtitle Banner
    tx_sub1 = slide1.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(3.0), Inches(0.45))
    p_tp = tx_sub1.text_frame.paragraphs[0]
    p_tp.text = "TITLE PAGE"
    p_tp.font.name = "Segoe UI"
    p_tp.font.size = Pt(18)
    p_tp.font.bold = True
    p_tp.font.color.rgb = BLACK

    # Left Column: SIH Metadata Fields
    tx_meta = slide1.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(6.8), Inches(4.1))
    tf_m = tx_meta.text_frame
    tf_m.word_wrap = True

    meta_rows = [
        ("• Problem Statement ID –", " SIH26197", DARK_GREEN, 16, True),
        ("• Problem Statement Title-", " Student Innovation—Ideas that showcase the rich cultural heritage and traditions of India.", DARK_GREEN, 12.5, False),
        ("• Theme-", " Heritage & Culture", DARK_GREEN, 14, False),
        ("• PS Category-", " Software", DARK_GREEN, 14, False),
        ("• Lead Organization-", " All India Council for Technical Education (AICTE)", DARK_GREEN, 13, False),
        ("• Team ID-", " [YOUR TEAM ID]", DARK_GREEN, 14, False),
        ("• Team Name-", " Team VirasatX", DARK_GREEN, 14, False)
    ]

    for idx, (lbl, val, col, sz, is_hl) in enumerate(meta_rows):
        p = tf_m.paragraphs[0] if idx == 0 else tf_m.add_paragraph()
        p.space_after = Pt(8)
        r1 = p.add_run()
        r1.text = lbl
        r1.font.name = "Segoe UI"
        r1.font.size = Pt(sz)
        r1.font.bold = True
        r1.font.color.rgb = BLACK

        r2 = p.add_run()
        r2.text = val
        r2.font.name = "Segoe UI"
        r2.font.size = Pt(sz + (2 if is_hl else 0))
        r2.font.bold = True
        r2.font.color.rgb = col

    # One-Line Project Pitch Card
    pitch_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.6), Inches(6.8), Inches(1.35))
    pitch_card.fill.solid()
    pitch_card.fill.fore_color.rgb = LIGHT_BG
    pitch_card.line.color.rgb = ROYAL_BLUE
    pitch_card.line.width = Pt(1.5)
    tf_pt = pitch_card.text_frame
    tf_pt.word_wrap = True
    p_pth = tf_pt.paragraphs[0]
    p_pth.text = "ONE-LINE PROJECT PITCH"
    p_pth.font.name = "Segoe UI"
    p_pth.font.size = Pt(10)
    p_pth.font.bold = True
    p_pth.font.color.rgb = ROYAL_BLUE
    p_pth.space_after = Pt(3)

    p_ptb = tf_pt.add_paragraph()
    p_ptb.text = "“An AI-powered digital heritage ecosystem that connects monuments, artifacts, traditions, places, people and knowledge into one discoverable cultural experience.”"
    p_ptb.font.name = "Segoe UI"
    p_ptb.font.size = Pt(12)
    p_ptb.font.italic = True
    p_ptb.font.color.rgb = DARK_GREY

    # Right Column: Prototype Hero Preview
    if os.path.exists(bulb_png):
        slide1.shapes.add_picture(bulb_png, Inches(8.5), Inches(1.1), width=Inches(3.8))
    
    if os.path.exists(home_img):
        hero_pic = slide1.shapes.add_picture(home_img, Inches(7.7), Inches(1.7), width=Inches(5.1))
        # Add border
        hero_pic.line.color.rgb = BORDER_GREY
        hero_pic.line.width = Pt(1.5)

    # Hero Tagline Pill
    tag_pill = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(6.45), Inches(4.3), Inches(0.45))
    tag_pill.fill.solid()
    tag_pill.fill.fore_color.rgb = NAVY
    tag_pill.line.fill.background()
    p_tg = tag_pill.text_frame.paragraphs[0]
    p_tg.text = "“India’s Heritage, Understood, Preserved and Experienced.”"
    p_tg.font.name = "Segoe UI"
    p_tg.font.size = Pt(10)
    p_tg.font.bold = True
    p_tg.font.color.rgb = WHITE
    p_tg.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 2: PROPOSED SOLUTION & ARCHITECTURE
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "VIRASATX", "(India's Heritage Repository)", is_pill=True, innovation_star=True)
    add_footer(slide2, 2)

    # Left Column Container Box 1: PROBLEM EXISTING
    prob_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.95), Inches(7.8), Inches(1.4))
    prob_box.fill.solid()
    prob_box.fill.fore_color.rgb = WHITE
    prob_box.line.color.rgb = ROYAL_BLUE
    prob_box.line.width = Pt(1.2)
    tf_pr = prob_box.text_frame
    tf_pr.word_wrap = True

    p = tf_pr.paragraphs[0]
    p.text = "PROBLEM EXISTING"
    p.font.name = "Segoe UI"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.space_after = Pt(2)

    prob_bullets = [
        ("➤ Fragmented Heritage Knowledge : ", "Monuments, archives, manuscripts, and living communities are scattered across disconnected portals."),
        ("➤ Passive Digital Discovery : ", "Existing portals present static text/listings without interconnected cultural context or visual guidance."),
        ("➤ Low Living Heritage Visibility : ", "Traditional craft guilds, master artisans, and oral traditions remain completely separated from monuments.")
    ]
    for b_title, b_desc in prob_bullets:
        p = tf_pr.add_paragraph()
        p.space_after = Pt(2)
        r_t = p.add_run()
        r_t.text = b_title
        r_t.font.name = "Segoe UI"
        r_t.font.size = Pt(9.5)
        r_t.font.bold = True
        r_t.font.color.rgb = ROYAL_BLUE

        r_d = p.add_run()
        r_d.text = b_desc
        r_d.font.name = "Segoe UI"
        r_d.font.size = Pt(9)
        r_d.font.color.rgb = BLACK

    # Middle 4 Screenshot Thumbnails Row
    thumbs = [
        (discover_img, "Discover Catalog", Inches(0.45)),
        (map_img, "Heritage Map", Inches(2.45)),
        (traditions_img, "Living Traditions", Inches(4.45)),
        (ai_img, "Virasat AI Guide", Inches(6.45))
    ]
    for img_path, lbl, x_pos in thumbs:
        if os.path.exists(img_path):
            pic = slide2.shapes.add_picture(img_path, x_pos, Inches(2.45), width=Inches(1.8), height=Inches(1.05))
            pic.line.color.rgb = BORDER_GREY
            pic.line.width = Pt(1)
            # Label box
            lbl_b = slide2.shapes.add_textbox(x_pos, Inches(3.52), Inches(1.8), Inches(0.25))
            p_l = lbl_b.text_frame.paragraphs[0]
            p_l.text = lbl
            p_l.font.name = "Segoe UI"
            p_l.font.size = Pt(8)
            p_l.font.bold = True
            p_l.font.color.rgb = DARK_GREY
            p_l.alignment = PP_ALIGN.CENTER

    # Left Box 2: PROPOSED SOLUTION
    sol_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(3.82), Inches(7.8), Inches(1.75))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = WHITE
    sol_box.line.color.rgb = ROYAL_BLUE
    sol_box.line.width = Pt(1.2)
    tf_sol = sol_box.text_frame
    tf_sol.word_wrap = True

    p = tf_sol.paragraphs[0]
    p.text = "PROPOSED SOLUTION"
    p.font.name = "Segoe UI"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.space_after = Pt(2)

    sol_bullets = [
        ("➤ Unified Heritage Repository : ", "Explore monuments, artifacts, collections, manuscripts, and living traditions in one place."),
        ("➤ Multimodal AI Guidance : ", "Ask Virasat AI for research inquiries and upload photos for visual iconography identification."),
        ("➤ Contextual Continuity : ", "Chronological 5,000-year timeline, geospatial GIS corridors, and ancient script paleography."),
        ("➤ Responsible Cultural Travel : ", "Carrying-capacity advisories and rural craft linkages to empower local artisan economies.")
    ]
    for b_title, b_desc in sol_bullets:
        p = tf_sol.add_paragraph()
        p.space_after = Pt(2)
        r_t = p.add_run()
        r_t.text = b_title
        r_t.font.name = "Segoe UI"
        r_t.font.size = Pt(9.5)
        r_t.font.bold = True
        r_t.font.color.rgb = ROYAL_BLUE

        r_d = p.add_run()
        r_d.text = b_desc
        r_d.font.name = "Segoe UI"
        r_d.font.size = Pt(9)
        r_d.font.color.rgb = BLACK

    # Left Box 3: UVP (UNIQUE VALUE PROPOSITION)
    uvp_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(5.65), Inches(7.8), Inches(1.3))
    uvp_box.fill.solid()
    uvp_box.fill.fore_color.rgb = LIGHT_BG
    uvp_box.line.color.rgb = BORDER_GREY
    uvp_box.line.width = Pt(1.2)
    tf_u = uvp_box.text_frame
    tf_u.word_wrap = True

    p = tf_u.paragraphs[0]
    p.text = "UVP (UNIQUE VALUE PROPOSITION)"
    p.font.name = "Segoe UI"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = SLATE
    p.space_after = Pt(2)

    p_flow = tf_u.add_paragraph()
    p_flow.text = "Artifact  ➔  Story  ➔  Period  ➔  Place  ➔  Tradition  ➔  Community  ➔  Learning  ➔  Responsible Experience"
    p_flow.font.name = "Segoe UI"
    p_flow.font.size = Pt(9.5)
    p_flow.font.bold = True
    p_flow.font.color.rgb = ROYAL_BLUE
    p_flow.space_after = Pt(2)

    p_qt = tf_u.add_paragraph()
    p_qt.text = "“VirasatX doesn't just display heritage — it connects the cultural context around it.”"
    p_qt.font.name = "Segoe UI"
    p_qt.font.size = Pt(11)
    p_qt.font.bold = True
    p_qt.font.italic = True
    p_qt.font.color.rgb = NAVY
    p_qt.alignment = PP_ALIGN.CENTER

    # Right Column: ARCHITECTURE Vertical Flowchart
    arch_outer = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.45), Inches(0.95), Inches(4.45), Inches(6.0))
    arch_outer.fill.solid()
    arch_outer.fill.fore_color.rgb = WHITE
    arch_outer.line.color.rgb = NAVY
    arch_outer.line.width = Pt(1.5)

    # Architecture Header Pill
    arch_pill = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.6), Inches(1.1), Inches(1.6), Inches(0.35))
    arch_pill.fill.solid()
    arch_pill.fill.fore_color.rgb = NAVY
    arch_pill.line.fill.background()
    p_ah = arch_pill.text_frame.paragraphs[0]
    p_ah.text = "ARCHITECTURE"
    p_ah.font.name = "Segoe UI"
    p_ah.font.size = Pt(10)
    p_ah.font.bold = True
    p_ah.font.color.rgb = WHITE
    p_ah.alignment = PP_ALIGN.CENTER

    # Node 1: User
    n1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.65), Inches(1.55), Inches(4.05), Inches(0.65))
    n1.fill.solid()
    n1.fill.fore_color.rgb = RGBColor(239, 246, 255)
    n1.line.color.rgb = RGBColor(147, 197, 253)
    p_n1 = n1.text_frame.paragraphs[0]
    p_n1.text = "👤 User / Patron / Student"
    p_n1.font.name = "Segoe UI"
    p_n1.font.size = Pt(10.5)
    p_n1.font.bold = True
    p_n1.font.color.rgb = RGBColor(30, 64, 175)
    p_n1.alignment = PP_ALIGN.CENTER
    p_n1b = n1.text_frame.add_paragraph()
    p_n1b.text = "Text Query • Photo Upload • GIS Map Exploration"
    p_n1b.font.name = "Segoe UI"
    p_n1b.font.size = Pt(8.5)
    p_n1b.font.color.rgb = DARK_GREY
    p_n1b.alignment = PP_ALIGN.CENTER

    # Arrow
    tx_a1 = slide2.shapes.add_textbox(Inches(10.5), Inches(2.2), Inches(0.4), Inches(0.25))
    p_arr1 = tx_a1.text_frame.paragraphs[0]
    p_arr1.text = "↓"
    p_arr1.font.size = Pt(12)
    p_arr1.font.bold = True
    p_arr1.font.color.rgb = ROYAL_BLUE
    p_arr1.alignment = PP_ALIGN.CENTER

    # Node 2: Core
    n2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.65), Inches(2.5), Inches(4.05), Inches(0.65))
    n2.fill.solid()
    n2.fill.fore_color.rgb = RGBColor(250, 248, 245)
    n2.line.color.rgb = GOLD
    p_n2 = n2.text_frame.paragraphs[0]
    p_n2.text = "🏛️ VirasatX Digital Heritage Core"
    p_n2.font.name = "Segoe UI"
    p_n2.font.size = Pt(10.5)
    p_n2.font.bold = True
    p_n2.font.color.rgb = GOLD
    p_n2.alignment = PP_ALIGN.CENTER
    p_n2b = n2.text_frame.add_paragraph()
    p_n2b.text = "Grounded Knowledge Catalog • Local Media Assets"
    p_n2b.font.name = "Segoe UI"
    p_n2b.font.size = Pt(8.5)
    p_n2b.font.color.rgb = DARK_GREY
    p_n2b.alignment = PP_ALIGN.CENTER

    # Arrow
    tx_a2 = slide2.shapes.add_textbox(Inches(10.5), Inches(3.15), Inches(0.4), Inches(0.25))
    p_arr2 = tx_a2.text_frame.paragraphs[0]
    p_arr2.text = "↓"
    p_arr2.font.size = Pt(12)
    p_arr2.font.bold = True
    p_arr2.font.color.rgb = ROYAL_BLUE
    p_arr2.alignment = PP_ALIGN.CENTER

    # 3 Mode Boxes Row
    modes = [
        ("ONLINE MODE\nInteractive 3D & Atlas", Inches(8.65), RGBColor(240, 249, 255), RGBColor(125, 211, 252), RGBColor(3, 105, 161)),
        ("AI RESEARCH\nGrounded AI & Vision", Inches(10.05), RGBColor(255, 247, 237), RGBColor(253, 186, 116), RGBColor(194, 65, 12)),
        ("LIVING HERITAGE\nGI Artisan Guilds", Inches(11.45), RGBColor(240, 253, 244), RGBColor(134, 239, 172), DARK_GREEN)
    ]
    for m_text, m_x, m_bg, m_border, m_col in modes:
        box_m = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, m_x, Inches(3.45), Inches(1.28), Inches(0.65))
        box_m.fill.solid()
        box_m.fill.fore_color.rgb = m_bg
        box_m.line.color.rgb = m_border
        p_m = box_m.text_frame.paragraphs[0]
        p_m.text = m_text
        p_m.font.name = "Segoe UI"
        p_m.font.size = Pt(8)
        p_m.font.bold = True
        p_m.font.color.rgb = m_col
        p_m.alignment = PP_ALIGN.CENTER

    # Arrow
    tx_a3 = slide2.shapes.add_textbox(Inches(10.5), Inches(4.15), Inches(0.4), Inches(0.25))
    p_arr3 = tx_a3.text_frame.paragraphs[0]
    p_arr3.text = "↓"
    p_arr3.font.size = Pt(12)
    p_arr3.font.bold = True
    p_arr3.font.color.rgb = ROYAL_BLUE
    p_arr3.alignment = PP_ALIGN.CENTER

    # Node 3: Grounded Evidence
    n3 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.65), Inches(4.45), Inches(4.05), Inches(0.65))
    n3.fill.solid()
    n3.fill.fore_color.rgb = RGBColor(240, 249, 255)
    n3.line.color.rgb = ROYAL_BLUE
    p_n3 = n3.text_frame.paragraphs[0]
    p_n3.text = "✔️ Grounded Evidence & Citations [ASI-104]"
    p_n3.font.name = "Segoe UI"
    p_n3.font.size = Pt(10)
    p_n3.font.bold = True
    p_n3.font.color.rgb = ROYAL_BLUE
    p_n3.alignment = PP_ALIGN.CENTER
    p_n3b = n3.text_frame.add_paragraph()
    p_n3b.text = "Confidence Ratings • Dynamic Curatorial Narration"
    p_n3b.font.name = "Segoe UI"
    p_n3b.font.size = Pt(8.5)
    p_n3b.font.color.rgb = DARK_GREY
    p_n3b.alignment = PP_ALIGN.CENTER

    # Arrow
    tx_a4 = slide2.shapes.add_textbox(Inches(10.5), Inches(5.15), Inches(0.4), Inches(0.25))
    p_arr4 = tx_a4.text_frame.paragraphs[0]
    p_arr4.text = "↓"
    p_arr4.font.size = Pt(12)
    p_arr4.font.bold = True
    p_arr4.font.color.rgb = ROYAL_BLUE
    p_arr4.alignment = PP_ALIGN.CENTER

    # Node 4: Unbroken Cultural Horizon Experience
    n4 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.65), Inches(5.45), Inches(4.05), Inches(0.65))
    n4.fill.solid()
    n4.fill.fore_color.rgb = NAVY
    n4.line.fill.background()
    p_n4 = n4.text_frame.paragraphs[0]
    p_n4.text = "✨ Unbroken Cultural Horizon Experience"
    p_n4.font.name = "Segoe UI"
    p_n4.font.size = Pt(10.5)
    p_n4.font.bold = True
    p_n4.font.color.rgb = WHITE
    p_n4.alignment = PP_ALIGN.CENTER
    p_n4b = n4.text_frame.add_paragraph()
    p_n4b.text = "Explore ➔ Learn ➔ Understand ➔ Preserve"
    p_n4b.font.name = "Segoe UI"
    p_n4b.font.size = Pt(8.5)
    p_n4b.font.color.rgb = RGBColor(186, 230, 253)
    p_n4b.alignment = PP_ALIGN.CENTER

    # Standards Row at Bottom
    std_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.65), Inches(6.25), Inches(4.05), Inches(0.45))
    std_box.fill.solid()
    std_box.fill.fore_color.rgb = LIGHT_BG
    std_box.line.color.rgb = BORDER_GREY
    p_std = std_box.text_frame.paragraphs[0]
    p_std.text = "🏛️ ASI Grounded   🔒 No PII   ♿ WCAG 2.1 AA   📜 Open Access"
    p_std.font.name = "Segoe UI"
    p_std.font.size = Pt(8.5)
    p_std.font.bold = True
    p_std.font.color.rgb = SLATE
    p_std.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "T E C H N I C A L   A P P R O A C H")
    add_footer(slide3, 3)

    # Top Exploration Flowchart Outer Container
    flow_outer = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.95), Inches(12.43), Inches(3.3))
    flow_outer.fill.solid()
    flow_outer.fill.fore_color.rgb = WHITE
    flow_outer.line.color.rgb = ROYAL_BLUE
    flow_outer.line.width = Pt(1.5)

    # Start Pill + Decision Box
    p_st = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.1), Inches(1.2), Inches(0.35))
    p_st.fill.solid()
    p_st.fill.fore_color.rgb = SLATE
    p_st.line.fill.background()
    p_st.text_frame.paragraphs[0].text = "🚀 Start Here"
    p_st.text_frame.paragraphs[0].font.size = Pt(9.5)
    p_st.text_frame.paragraphs[0].font.bold = True
    p_st.text_frame.paragraphs[0].font.color.rgb = WHITE
    p_st.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    p_dec = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.9), Inches(1.1), Inches(3.2), Inches(0.35))
    p_dec.fill.solid()
    p_dec.fill.fore_color.rgb = DARK_GREY
    p_dec.line.fill.background()
    p_dec.text_frame.paragraphs[0].text = "Choose Cultural Exploration Pathway"
    p_dec.text_frame.paragraphs[0].font.size = Pt(10)
    p_dec.text_frame.paragraphs[0].font.bold = True
    p_dec.text_frame.paragraphs[0].font.color.rgb = WHITE
    p_dec.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Pathway Labels
    tx_pw = slide3.shapes.add_textbox(Inches(5.3), Inches(1.1), Inches(7.4), Inches(0.35))
    tf_pw = tx_pw.text_frame
    p_pw = tf_pw.paragraphs[0]
    r = p_pw.add_run(); r.text = "• Visual Specimen Scan   "; r.font.bold = True; r.font.color.rgb = ORANGE; r.font.size = Pt(9.5)
    r = p_pw.add_run(); r.text = "• Scholarly Inquiry / Text   "; r.font.bold = True; r.font.color.rgb = ROYAL_BLUE; r.font.size = Pt(9.5)
    r = p_pw.add_run(); r.text = "• Geospatial & 3D Studio"; r.font.bold = True; r.font.color.rgb = DARK_GREEN; r.font.size = Pt(9.5)

    # 3 Pipeline Columns
    cols_data = [
        ("VISION AI PATHWAY", ORANGE, RGBColor(254, 215, 170), Inches(0.6), [
            "📷 Upload Artifact / Monument Photo",
            "⚙️ Base64 Preprocessing & Client Sanitization",
            "🤖 Gemini 2.5 Flash Vision Multimodal Engine",
            "🔍 Identify Mudras, Asanas, Drapery & Dynasty",
            "🎯 Output: Classification + Confidence Meter"
        ]),
        ("AI HERITAGE GUIDE (RAG)", ROYAL_BLUE, RGBColor(191, 219, 254), Inches(4.7), [
            "💬 Natural Query (English, Hindi, Sanskrit)",
            "📚 Lexical & Semantic Context Vector Retrieval",
            "🏛️ Ground against Curated ASI Accession Records",
            "⚡ Serverless Backend API Execution (/api/ai-guide)",
            "🎯 Output: Grounded Response with [ASI-104] Citations"
        ]),
        ("GEOSPATIAL & 3D STUDIO", DARK_GREEN, RGBColor(187, 247, 208), Inches(8.8), [
            "🗺️ Leaflet OSM Interactive GPS Coordinate Atlas",
            "🧭 6 Thematic Cultural Corridors (Chola, etc.)",
            "🎮 Three.js WebGL 360° Tactile Inspection Studio",
            "💡 Raking, Amber & Neutral Museum Lighting",
            "🎯 Output: Tactile Inspection & Audio Narration"
        ])
    ]

    for title, header_bg, border_col, col_x, steps in cols_data:
        # Col box
        c_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_x, Inches(1.55), Inches(3.9), Inches(2.55))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = LIGHT_BG
        c_box.line.color.rgb = border_col
        c_box.line.width = Pt(1.2)

        # Header
        h_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_x + Inches(0.1), Inches(1.65), Inches(3.7), Inches(0.35))
        h_box.fill.solid()
        h_box.fill.fore_color.rgb = header_bg
        h_box.line.fill.background()
        p_h = h_box.text_frame.paragraphs[0]
        p_h.text = title
        p_h.font.name = "Segoe UI"
        p_h.font.size = Pt(10)
        p_h.font.bold = True
        p_h.font.color.rgb = WHITE
        p_h.alignment = PP_ALIGN.CENTER

        # Steps
        for s_idx, step_txt in enumerate(steps):
            s_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_x + Inches(0.1), Inches(2.08 + s_idx * 0.38), Inches(3.7), Inches(0.32))
            s_box.fill.solid()
            s_box.fill.fore_color.rgb = WHITE
            s_box.line.color.rgb = border_col
            p_s = s_box.text_frame.paragraphs[0]
            p_s.text = step_txt
            p_s.font.name = "Segoe UI"
            p_s.font.size = Pt(8.5)
            p_s.font.bold = (s_idx == 4)
            p_s.font.color.rgb = header_bg if (s_idx == 4) else BLACK

    # Bottom Row: Tech Stack Grid & Architecture Callout
    ts_outer = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(4.38), Inches(7.4), Inches(2.6))
    ts_outer.fill.solid()
    ts_outer.fill.fore_color.rgb = WHITE
    ts_outer.line.color.rgb = ROYAL_BLUE
    ts_outer.line.width = Pt(1.2)

    # Tech Stack Pill
    ts_pill = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.5), Inches(1.5), Inches(0.3))
    ts_pill.fill.solid()
    ts_pill.fill.fore_color.rgb = NAVY
    ts_pill.line.fill.background()
    p_tsp = ts_pill.text_frame.paragraphs[0]
    p_tsp.text = "TECH STACK"
    p_tsp.font.name = "Segoe UI"
    p_tsp.font.size = Pt(9.5)
    p_tsp.font.bold = True
    p_tsp.font.color.rgb = WHITE
    p_tsp.alignment = PP_ALIGN.CENTER

    tech_items = [
        ("Frontend", "React 19 • Vite 6", Inches(0.6), Inches(4.9)),
        ("Language & CSS", "TypeScript • Tailwind v4", Inches(2.4), Inches(4.9)),
        ("AI & Multimodal", "Gemini 2.5 Flash", Inches(4.2), Inches(4.9)),
        ("AI SDK", "@google/genai", Inches(6.0), Inches(4.9)),
        ("3D Engine", "Three.js (WebGL)", Inches(0.6), Inches(5.8)),
        ("Mapping GIS", "Leaflet • OSM", Inches(2.4), Inches(5.8)),
        ("Backend & DB", "Supabase • PostgreSQL", Inches(4.2), Inches(5.8)),
        ("Deployment", "Vercel Serverless", Inches(6.0), Inches(5.8))
    ]
    for cat, val, tx, ty in tech_items:
        t_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx, ty, Inches(1.7), Inches(0.75))
        t_box.fill.solid()
        t_box.fill.fore_color.rgb = LIGHT_BG
        t_box.line.color.rgb = BORDER_GREY
        tf_t = t_box.text_frame
        p1 = tf_t.paragraphs[0]
        p1.text = cat
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(8)
        p1.font.bold = True
        p1.font.color.rgb = ROYAL_BLUE
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf_t.add_paragraph()
        p2.text = val
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(8.5)
        p2.font.bold = True
        p2.font.color.rgb = BLACK
        p2.alignment = PP_ALIGN.CENTER

    # Prototype Verification Callout Box
    pv_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.05), Inches(4.38), Inches(4.83), Inches(2.6))
    pv_box.fill.solid()
    pv_box.fill.fore_color.rgb = LIGHT_BG
    pv_box.line.color.rgb = BORDER_GREY
    pv_box.line.width = Pt(1.2)
    tf_pv = pv_box.text_frame
    tf_pv.word_wrap = True

    p = tf_pv.paragraphs[0]
    p.text = "PROTOTYPE VERIFICATION"
    p.font.name = "Segoe UI"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(4)

    pv_bullets = [
        ("➤ Secret Shielding : ", "Gemini API keys execute strictly in backend serverless functions; zero keys in client bundles."),
        ("➤ Defensive SafeImage : ", "Shimmer skeletons with zero layout shift (CLS < 0.05) and archival parchment fallback.")
    ]
    for b_t, b_d in pv_bullets:
        p = tf_pv.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run(); r.text = b_t; r.font.bold = True; r.font.size = Pt(9.5); r.font.color.rgb = ROYAL_BLUE
        r = p.add_run(); r.text = b_d; r.font.size = Pt(9); r.font.color.rgb = BLACK

    # Principle Pill
    pr_pill = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(6.35), Inches(4.5), Inches(0.42))
    pr_pill.fill.solid()
    pr_pill.fill.fore_color.rgb = NAVY
    pr_pill.line.fill.background()
    p_pr = pr_pill.text_frame.paragraphs[0]
    p_pr.text = "Principle : Modular • Scalable • Secure • Evidence-Aware"
    p_pr.font.name = "Segoe UI"
    p_pr.font.size = Pt(9.5)
    p_pr.font.bold = True
    p_pr.font.color.rgb = WHITE
    p_pr.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "F E A S I B I L I T Y   A N D   V I A B I L I T Y")
    add_footer(slide4, 4)

    # Top Left: FEASIBILITY
    feas_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.95), Inches(6.05), Inches(3.5))
    feas_box.fill.solid()
    feas_box.fill.fore_color.rgb = WHITE
    feas_box.line.color.rgb = ROYAL_BLUE
    feas_box.line.width = Pt(1.5)
    tf_f = feas_box.text_frame
    tf_f.word_wrap = True

    p = tf_f.paragraphs[0]
    p.text = "FEASIBILITY :"
    p.font.name = "Segoe UI"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.space_after = Pt(4)

    feas_bullets = [
        ("Technical : ", "Proven Tech Stack — React 19, TypeScript, Vite, Supabase PostgreSQL, and Three.js. Smooth cross-platform web execution with zero heavy native app download requirements."),
        ("Secret Protection & Security : ", "Serverless execution ensures Gemini API keys and database credentials are strictly isolated from client-side bundles."),
        ("Innovation : ", "Domain-grounded RAG linking inquiries directly to verified ASI accession records (e.g. [ASI-104]) with confidence scoring."),
        ("Operational : ", "Runs efficiently on low-bandwidth connections through client-side tile caching and progressive <SafeImage /> fallbacks."),
        ("Cultural Data Standards : ", "Uses open cultural access datasets structured to ASI and National Mission for Manuscripts (NMM) archival standards.")
    ]
    for b_t, b_d in feas_bullets:
        p = tf_f.add_paragraph()
        p.space_after = Pt(3)
        r = p.add_run(); r.text = b_t; r.font.bold = True; r.font.size = Pt(9); r.font.color.rgb = BLACK
        r = p.add_run(); r.text = b_d; r.font.size = Pt(8.5); r.font.color.rgb = DARK_GREY

    # Top Right: VIABILITY
    viab_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(0.95), Inches(6.05), Inches(3.5))
    viab_box.fill.solid()
    viab_box.fill.fore_color.rgb = WHITE
    viab_box.line.color.rgb = DARK_GREEN
    viab_box.line.width = Pt(1.5)
    tf_v = viab_box.text_frame
    tf_v.word_wrap = True

    p = tf_v.paragraphs[0]
    p.text = "VIABILITY :"
    p.font.name = "Segoe UI"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.space_after = Pt(4)

    viab_bullets = [
        ("Educational Reach : ", "Enables schools, rural colleges, and researchers across India to inspect high-resolution monuments, manuscripts, and 3D bronzes remotely."),
        ("Policy & National Mission : ", "Directly supports AICTE Student Innovation, Ministry of Culture digitization drives, and UNESCO Sustainable Tourism (SDG 11)."),
        ("Infrastructure Overhead : ", "Serverless edge deployment and Supabase free-tier database keep initial operational costs near zero, scaling seamlessly with traffic."),
        ("Living Heritage Inclusion : ", "Directly spotlights GI-certified master artisan cooperatives (Swamimalai Bronzes, Aranmula Mirrors, Patan Patola) to foster rural craft livelihoods."),
        ("Scalability & Expansion : ", "Modular database schema allows continuous ingestion of state archaeological records and regional manuscript repositories.")
    ]
    for b_t, b_d in viab_bullets:
        p = tf_v.add_paragraph()
        p.space_after = Pt(3)
        r = p.add_run(); r.text = b_t; r.font.bold = True; r.font.size = Pt(9); r.font.color.rgb = BLACK
        r = p.add_run(); r.text = b_d; r.font.size = Pt(8.5); r.font.color.rgb = DARK_GREY

    # Bottom Left: TECHNICAL CHALLENGES
    tc_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(4.6), Inches(6.05), Inches(2.35))
    tc_box.fill.solid()
    tc_box.fill.fore_color.rgb = LIGHT_BG
    tc_box.line.color.rgb = BORDER_GREY
    tc_box.line.width = Pt(1.2)
    tf_tc = tc_box.text_frame
    tf_tc.word_wrap = True

    p = tf_tc.paragraphs[0]
    p.text = "TECHNICAL CHALLENGES :"
    p.font.name = "Segoe UI"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.space_after = Pt(4)

    p = tf_tc.add_paragraph()
    p.space_after = Pt(3)
    r = p.add_run(); r.text = "Risk : "; r.font.bold = True; r.font.size = Pt(9.5); r.font.color.rgb = BLACK
    r = p.add_run(); r.text = "Generative AI models risk hallucinating historical dates, dynastic chronologies, and mudra attributions."; r.font.size = Pt(9); r.font.color.rgb = DARK_GREY

    p = tf_tc.add_paragraph()
    r = p.add_run(); r.text = "Mitigation : "; r.font.bold = True; r.font.size = Pt(9.5); r.font.color.rgb = BLACK
    r = p.add_run(); r.text = "Inquiries are strictly grounded against curated ASI accession catalogs; prompts enforce low temperature (0.2), cite explicit record numbers, and display transparent confidence meters."; r.font.size = Pt(9); r.font.color.rgb = DARK_GREY

    # Bottom Right: CULTURAL & ADOPTION CHALLENGES
    cc_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.6), Inches(6.05), Inches(2.35))
    cc_box.fill.solid()
    cc_box.fill.fore_color.rgb = LIGHT_BG
    cc_box.line.color.rgb = BORDER_GREY
    cc_box.line.width = Pt(1.2)
    tf_cc = cc_box.text_frame
    tf_cc.word_wrap = True

    p = tf_cc.paragraphs[0]
    p.text = "CULTURAL & ADOPTION CHALLENGES :"
    p.font.name = "Segoe UI"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.space_after = Pt(4)

    p = tf_cc.add_paragraph()
    p.space_after = Pt(3)
    r = p.add_run(); r.text = "Risk : "; r.font.bold = True; r.font.size = Pt(9.5); r.font.color.rgb = BLACK
    r = p.add_run(); r.text = "Traditional artisans and master sthapatis may face privacy intrusion or contact spam from public directories."; r.font.size = Pt(9); r.font.color.rgb = DARK_GREY

    p = tf_cc.add_paragraph()
    r = p.add_run(); r.text = "Mitigation : "; r.font.bold = True; r.font.size = Pt(9.5); r.font.color.rgb = BLACK
    r = p.add_run(); r.text = "Ethical custodianship protocol connects patrons strictly with certified Geographical Indication (GI) cooperative offices with zero exposure of personal telephone numbers or private home addresses."; r.font.size = Pt(9); r.font.color.rgb = DARK_GREY

    # =========================================================================
    # SLIDE 5: IMPACTS & BENEFITS
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "I M P A C T S   &   B E N E F I T S")
    add_footer(slide5, 5)

    # Top Left: 3 Core Impact Points
    tx_imp = slide5.shapes.add_textbox(Inches(0.45), Inches(0.95), Inches(7.4), Inches(2.75))
    tf_im = tx_imp.text_frame
    tf_im.word_wrap = True

    imp_bullets = [
        ("• Cultural Preservation & Digitization : ", "Digitally structures and safeguards fragile palm-leaf and birch-bark manuscripts, epigraphical records, and vanishing oral craft knowledge before physical decay occurs."),
        ("• Educational & Youth Engagement : ", "Replaces passive textbook memorization with tactile 360° WebGL inspection, interactive geospatial timelines, and voice-assisted exploration for schools and universities nationwide."),
        ("• Artisan Livelihood & Sustainable Tourism (SDG 11) : ", "Connects cultural tourists with GI-certified master artisan guilds and disperses tourist crowds from overburdened monuments to lesser-known heritage clusters.")
    ]
    for idx, (b_t, b_d) in enumerate(imp_bullets):
        p = tf_im.paragraphs[0] if idx == 0 else tf_im.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run(); r.text = b_t; r.font.name = "Segoe UI"; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = BLACK
        r = p.add_run(); r.text = b_d; r.font.name = "Segoe UI"; r.font.size = Pt(10); r.font.color.rgb = DARK_GREY

    # Top Right: Feature Comparison Table
    table_shape = slide5.shapes.add_table(9, 3, Inches(8.15), Inches(0.95), Inches(4.7), Inches(2.75))
    tbl = table_shape.table
    tbl.columns[0].width = Inches(2.7)
    tbl.columns[1].width = Inches(1.0)
    tbl.columns[2].width = Inches(1.0)

    features = [
        ("Feature", "Existing Portals", "VirasatX"),
        ("Interconnected Cultural Context", "❌", "✔️"),
        ("360° WebGL Archival Studio", "❌", "✔️"),
        ("Visual Iconography AI Scanner", "❌", "✔️"),
        ("Grounded ASI Accession Citations", "❌", "✔️"),
        ("Living GI Artisan Guild Linkages", "❌", "✔️"),
        ("Ancient Manuscript Decipherment", "❌", "✔️"),
        ("Responsible Travel & Carrying Capacity", "❌", "✔️"),
        ("100% Local Media SafeImage Pipeline", "❌", "✔️")
    ]

    for row_idx, row_data in enumerate(features):
        for col_idx, cell_value in enumerate(row_data):
            cell = tbl.cell(row_idx, col_idx)
            cell.text = cell_value
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            p.font.name = "Segoe UI"
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                p.font.bold = True
                p.font.size = Pt(9)
                p.font.color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG if row_idx % 2 == 1 else WHITE
                p.font.size = Pt(8.5)
                p.font.bold = True
                if col_idx == 1:
                    p.font.color.rgb = RGBColor(220, 38, 38)
                elif col_idx == 2:
                    p.font.color.rgb = DARK_GREEN
                else:
                    p.font.color.rgb = BLACK

    # Middle Banner: HOW VIRASATX WORKS ACROSS MODES
    mid_bar = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(3.85), Inches(12.4), Inches(0.35))
    mid_bar.fill.solid()
    mid_bar.fill.fore_color.rgb = LIGHT_BG
    mid_bar.line.color.rgb = BORDER_GREY
    p_mb = mid_bar.text_frame.paragraphs[0]
    p_mb.text = "H O W   V I R A S A T X   W O R K S   A C R O S S   M O D E S"
    p_mb.font.name = "Segoe UI"
    p_mb.font.size = Pt(10)
    p_mb.font.bold = True
    p_mb.font.color.rgb = NAVY
    p_mb.alignment = PP_ALIGN.CENTER

    # Bottom 3 Colored Mode Cards
    modes_data = [
        ("MODE 1: EXPLORER (Public & Students)", ROYAL_BLUE, Inches(0.45), [
            ("• For : ", "Students, tourists, and casual cultural enthusiasts."),
            ("• Process : ", "Interactive search ➔ 3D model inspection ➔ Audio narrative ➔ Timeline navigation."),
            ("• Time : ", "Instant (sub-second query response)."),
            ("• Success Rate : ", "100% curated client availability."),
            ("• Best For : ", "Visual learning, monument discovery, school projects.")
        ]),
        ("MODE 2: SCHOLAR & RESEARCHER", ORANGE, Inches(4.7), [
            ("• For : ", "Historians, archaeologists, and epigraphists."),
            ("• Process : ", "Ask Virasat AI ➔ Vector retrieval ➔ Accession citations ➔ Script transcription."),
            ("• Time : ", "1-3 seconds RAG inference."),
            ("• Success Rate : ", "95%+ citation grounding accuracy."),
            ("• Best For : ", "Iconographic analysis, academic research, catalog citations.")
        ]),
        ("MODE 3: LIVING HERITAGE & TRAVEL", DARK_GREEN, Inches(8.95), [
            ("• For : ", "Cultural travelers and craft patrons."),
            ("• Process : ", "Plan Visit ➔ Dawn/dusk crowd advisories ➔ Discover nearby GI crafts ➔ Cooperative routing."),
            ("• Time : ", "Real-time spatial corridor calculation."),
            ("• Success Rate : ", "100% verified GI cooperative registry."),
            ("• Best For : ", "Ethical cultural itineraries, craft preservation, rural economy.")
        ])
    ]

    for m_title, m_col, m_x, m_bullets in modes_data:
        # Outer
        m_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, m_x, Inches(4.3), Inches(3.9), Inches(2.7))
        m_box.fill.solid()
        m_box.fill.fore_color.rgb = WHITE
        m_box.line.color.rgb = m_col
        m_box.line.width = Pt(1.5)

        # Header
        mh_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, m_x, Inches(4.3), Inches(3.9), Inches(0.42))
        mh_box.fill.solid()
        mh_box.fill.fore_color.rgb = m_col
        mh_box.line.fill.background()
        p_mh = mh_box.text_frame.paragraphs[0]
        p_mh.text = m_title
        p_mh.font.name = "Segoe UI"
        p_mh.font.size = Pt(9.5)
        p_mh.font.bold = True
        p_mh.font.color.rgb = WHITE
        p_mh.alignment = PP_ALIGN.CENTER

        # Text
        tx_mc = slide5.shapes.add_textbox(m_x + Inches(0.1), Inches(4.78), Inches(3.7), Inches(2.15))
        tf_mc = tx_mc.text_frame
        tf_mc.word_wrap = True
        for b_idx, (b_lbl, b_val) in enumerate(m_bullets):
            p = tf_mc.paragraphs[0] if b_idx == 0 else tf_mc.add_paragraph()
            p.space_after = Pt(2.5)
            r = p.add_run(); r.text = b_lbl; r.font.bold = True; r.font.size = Pt(8.5); r.font.color.rgb = BLACK
            r = p.add_run(); r.text = b_val; r.font.size = Pt(8.5); r.font.color.rgb = DARK_GREY

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "R E S E A R C H   A N D   R E F E R E N C E S")
    add_footer(slide6, 6)

    # Top Box: Official References & Standards
    ref_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(0.95), Inches(12.4), Inches(1.8))
    ref_box.fill.solid()
    ref_box.fill.fore_color.rgb = LIGHT_BG
    ref_box.line.color.rgb = BORDER_GREY
    ref_box.line.width = Pt(1.2)
    tf_rf = ref_box.text_frame
    tf_rf.word_wrap = True

    refs = [
        ("• Archaeological Survey of India (ASI) : ", "Primary repository of centrally protected monuments, excavation records, and accession catalogs. Links : asi.nic.in / NMMA"),
        ("• National Mission for Manuscripts (NMM) : ", "National database of ancient Indian manuscripts, scripts, and preservation guidelines. Links : namami.gov.in / IGNCA"),
        ("• Geographical Indications Registry of India : ", "Official documentation for GI-certified traditional craft practices (GI-029 Swamimalai, GI-007 Aranmula). Link : ipindia.gov.in"),
        ("• UNESCO World Heritage Centre : ", "Cultural criteria and sustainable tourism framework for monumental conservation (SDG 11). Link : whc.unesco.org")
    ]
    for idx, (r_lbl, r_val) in enumerate(refs):
        p = tf_rf.paragraphs[0] if idx == 0 else tf_rf.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run(); r.text = r_lbl; r.font.bold = True; r.font.size = Pt(9.5); r.font.color.rgb = BLACK
        r = p.add_run(); r.text = r_val; r.font.size = Pt(9); r.font.color.rgb = DARK_GREY

    # Middle Comparison: SOLUTIONS ALREADY EXIST ➔ OUR SOLUTION STANDS OUT
    # Left: Existing
    ex_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(2.9), Inches(5.7), Inches(1.75))
    ex_box.fill.solid()
    ex_box.fill.fore_color.rgb = RGBColor(240, 249, 255)
    ex_box.line.color.rgb = RGBColor(186, 230, 253)
    tf_ex = ex_box.text_frame
    tf_ex.word_wrap = True

    p = tf_ex.paragraphs[0]
    p.text = "SOLUTIONS ALREADY EXIST"
    p.font.name = "Segoe UI"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.space_after = Pt(3)

    p = tf_ex.add_paragraph()
    p.space_after = Pt(2)
    p.text = "• Most existing portals are isolated, static, and text-heavy; designed only for administrative filing."
    p.font.size = Pt(9); p.font.color.rgb = DARK_GREY

    p = tf_ex.add_paragraph()
    p.space_after = Pt(2)
    p.text = "• Separate websites for monuments, museums, manuscripts, and tourism with zero cross-linking."
    p.font.size = Pt(9); p.font.color.rgb = DARK_GREY

    p = tf_ex.add_paragraph()
    r = p.add_run(); r.text = "• AVAILABLE SITES : "; r.font.bold = True; r.font.size = Pt(9); r.font.color.rgb = BLACK
    r = p.add_run(); r.text = "ASI Portal / NMMA / State Tourism Portals / Commercial Travel Apps"; r.font.size = Pt(9); r.font.color.rgb = ROYAL_BLUE

    # Center Arrow
    tx_ar = slide6.shapes.add_textbox(Inches(6.3), Inches(3.5), Inches(0.7), Inches(0.5))
    p_ar = tx_ar.text_frame.paragraphs[0]
    p_ar.text = "➔"
    p_ar.font.size = Pt(26)
    p_ar.font.bold = True
    p_ar.font.color.rgb = BLACK
    p_ar.alignment = PP_ALIGN.CENTER

    # Right: Our Solution
    out_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.15), Inches(2.9), Inches(5.7), Inches(1.75))
    out_box.fill.solid()
    out_box.fill.fore_color.rgb = WHITE
    out_box.line.color.rgb = ROYAL_BLUE
    out_box.line.width = Pt(1.5)
    tf_out = out_box.text_frame
    tf_out.word_wrap = True

    p = tf_out.paragraphs[0]
    p.text = "OUR SOLUTION STANDS OUT"
    p.font.name = "Segoe UI"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.space_after = Pt(3)

    out_bullets = [
        ("• One Ecosystem : ", "Unites monuments, artifacts, manuscripts, and living traditions."),
        ("• Evidence-Aware : ", "Real-time RAG citing exact ASI accession records [ASI-104]."),
        ("• Tactile & Geospatial : ", "360° Three.js WebGL studio and 6-corridor Leaflet atlas."),
        ("• Ethical Custodianship : ", "Promotes GI craft guilds with zero personal contact exposure.")
    ]
    for b_l, b_v in out_bullets:
        p = tf_out.add_paragraph()
        p.space_after = Pt(1.5)
        r = p.add_run(); r.text = b_l; r.font.bold = True; r.font.size = Pt(9); r.font.color.rgb = BLACK
        r = p.add_run(); r.text = b_v; r.font.size = Pt(9); r.font.color.rgb = DARK_GREY

    # Bottom Row: Project Resources & UI Screenshots
    res_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(4.8), Inches(6.0), Inches(2.2))
    res_box.fill.solid()
    res_box.fill.fore_color.rgb = LIGHT_BG
    res_box.line.color.rgb = BORDER_GREY
    res_box.line.width = Pt(1.2)
    tf_rb = res_box.text_frame
    tf_rb.word_wrap = True

    p = tf_rb.paragraphs[0]
    p.text = "PROJECT RESOURCES"
    p.font.name = "Segoe UI"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ROYAL_BLUE
    p.space_after = Pt(4)

    res_items = [
        ("Live Prototype : ", "https://virasatxai.vercel.app/"),
        ("GitHub Repository : ", "https://github.com/ShrilCarpenter/VirasatX"),
        ("Full Documentation : ", "Architecture, API Specs & License in README.md"),
        ("Problem Statement : ", "SIH26197 • Heritage & Culture • AICTE")
    ]
    for r_l, r_v in res_items:
        p = tf_rb.add_paragraph()
        p.space_after = Pt(3)
        r = p.add_run(); r.text = r_l; r.font.bold = True; r.font.size = Pt(9.5); r.font.color.rgb = BLACK
        r = p.add_run(); r.text = r_v; r.font.size = Pt(9); r.font.color.rgb = ROYAL_BLUE if "http" in r_v else DARK_GREY

    # Bottom Right: 2 Screenshots
    if os.path.exists(home_img):
        pic1 = slide6.shapes.add_picture(home_img, Inches(6.65), Inches(4.8), width=Inches(3.0), height=Inches(1.8))
        pic1.line.color.rgb = BORDER_GREY
        pic1.line.width = Pt(1)
        # Caption
        tx_c1 = slide6.shapes.add_textbox(Inches(6.65), Inches(6.62), Inches(3.0), Inches(0.3))
        p_c1 = tx_c1.text_frame.paragraphs[0]
        p_c1.text = "Interactive Platform"
        p_c1.font.size = Pt(8.5); p_c1.font.bold = True; p_c1.font.color.rgb = DARK_GREY; p_c1.alignment = PP_ALIGN.CENTER

    if os.path.exists(ai_img):
        pic2 = slide6.shapes.add_picture(ai_img, Inches(9.85), Inches(4.8), width=Inches(3.0), height=Inches(1.8))
        pic2.line.color.rgb = BORDER_GREY
        pic2.line.width = Pt(1)
        # Caption
        tx_c2 = slide6.shapes.add_textbox(Inches(9.85), Inches(6.62), Inches(3.0), Inches(0.3))
        p_c2 = tx_c2.text_frame.paragraphs[0]
        p_c2.text = "Virasat AI Guide"
        p_c2.font.size = Pt(8.5); p_c2.font.bold = True; p_c2.font.color.rgb = DARK_GREY; p_c2.alignment = PP_ALIGN.CENTER

    # Save presentation
    prs.save(output_pptx_path)
    print(f"Successfully generated PowerPoint presentation at: {output_pptx_path}")

if __name__ == "__main__":
    out_dir = os.path.dirname(os.path.abspath(__file__))
    target_pptx = os.path.join(out_dir, "SIH2026_VirasatX_Idea_Submission.pptx")
    build_presentation(target_pptx)
